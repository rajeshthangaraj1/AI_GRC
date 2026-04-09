"""
data_embedding_ppt.py
Extracts text from ISO 42001 PowerPoint files and upserts into the existing
grc_docs Qdrant collection (alongside the Excel-embedded data).

Run AFTER data_embedding_v2.py:
    source .venv/bin/activate
    python data_embedding_ppt.py
"""

import os
import uuid
import re
from pptx import Presentation
from sentence_transformers import SentenceTransformer
from qdrant_client import QdrantClient
from qdrant_client.models import PointStruct

# =========================
# CONFIG
# =========================
PPTX_DIR    = "/home/rmrobot/Desktop/Rajesh/AI_GRC/GRC_Requirement"
MODEL_PATH  = "/home/rmrobot/Desktop/Rajesh/AI_GRC/models/bge-large"
COLLECTION  = "grc_docs"
QDRANT_PATH = "./data"
BATCH_SIZE  = 50
FRAMEWORK   = "ISO"
MIN_CHARS   = 30          # skip slides with less content than this

# PPT files to process (in GRC_Requirement/)
PPTX_FILES = [
    "AIMS-IMP-03-Implementation-Guide-for-AI-Management-Systems.pptx",
    "AIMS-IMP-04-The-ISO-Certification-Process-Your-Complete-Guide.pptx",
    "AIMS-IMP-05-Integrated-Management-Systems-Streamlining-Standards.pptx",
    "AIMS-PPT-01-ISO42001-Training - Year 1.pptx",
    "AIMS-PPT-02-ISO42001-Implementing-ISO-42001.pptx",
]

# =========================
# HELPERS
# =========================
def clean_text(text):
    text = re.sub(r"\s+", " ", str(text))
    return text.strip()

def extract_slide_text(slide):
    """Return (title_text, body_text) from a slide."""
    title = ""
    body_parts = []

    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        raw = " ".join(
            p.text.strip()
            for p in shape.text_frame.paragraphs
            if p.text.strip()
        )
        if not raw:
            continue
        # Detect title placeholder
        if shape.shape_type == 13:   # MSO_SHAPE_TYPE.PICTURE — skip
            continue
        try:
            if shape.is_placeholder and shape.placeholder_format.idx == 0:
                title = raw
                continue
        except Exception:
            pass
        body_parts.append(raw)

    return clean_text(title), clean_text(" | ".join(body_parts))

# =========================
# LOAD MODEL & QDRANT
# =========================
print("Loading embedding model...")
model = SentenceTransformer(MODEL_PATH, device="cuda")

client = QdrantClient(path=QDRANT_PATH)

if not client.collection_exists(COLLECTION):
    raise RuntimeError(
        f"Collection '{COLLECTION}' not found. "
        "Run data_embedding.py first to create it."
    )

# =========================
# PROCESS PPT FILES
# =========================
total_upserted = 0

for filename in PPTX_FILES:
    filepath = os.path.join(PPTX_DIR, filename)
    if not os.path.exists(filepath):
        print(f"[SKIP] Not found: {filename}")
        continue

    # Sheet label = short name without extension, e.g. "AIMS-IMP-03"
    sheet_label = filename.split("-")[0] + "-" + filename.split("-")[1] + "-" + filename.split("-")[2]
    sheet_label = os.path.splitext(filename)[0]   # full name without .pptx

    print(f"\nProcessing: {filename}")
    prs = Presentation(filepath)

    points = []
    slide_count = 0
    skipped = 0

    for slide_num, slide in enumerate(prs.slides, start=1):
        title, body = extract_slide_text(slide)

        combined = f"{title} {body}".strip()
        if len(combined) < MIN_CHARS:
            skipped += 1
            continue

        # Full embedding text — rich context prefix
        text = clean_text(
            f"This is ISO 42001 implementation content from {filename}. "
            f"Slide {slide_num}: {title}. Content: {body}"
        )

        # BM25 text — no filename prefix, lower-cased for keyword matching
        bm25_text = clean_text(f"Slide {slide_num}: {title}. {body}").lower()

        embedding = model.encode(
            f"passage: {text}",
            normalize_embeddings=True
        )

        points.append(PointStruct(
            id=str(uuid.uuid4()),
            vector=embedding.tolist(),
            payload={
                "text":       text,
                "bm25_text":  bm25_text,
                "framework":  FRAMEWORK,
                "sheet":      sheet_label,
            },
        ))
        slide_count += 1

        if len(points) >= BATCH_SIZE:
            client.upsert(collection_name=COLLECTION, points=points)
            total_upserted += len(points)
            points = []

    if points:
        client.upsert(collection_name=COLLECTION, points=points)
        total_upserted += len(points)
        points = []

    print(f"  Slides embedded: {slide_count} | Skipped (blank): {skipped}")

# =========================
# VERIFY FINAL COUNT
# =========================
scroll = client.scroll(collection_name=COLLECTION, limit=20000)
all_points = scroll[0]

framework_dist = {}
for p in all_points:
    fw = p.payload.get("framework", "Unknown")
    framework_dist[fw] = framework_dist.get(fw, 0) + 1

print(f"\nDONE: {total_upserted} PPT slides upserted into '{COLLECTION}'")
print(f"Total collection size: {len(all_points)} points")
print("\nFramework Distribution:")
for k, v in framework_dist.items():
    print(f"  {k}: {v}")

client.close()
